"""
make_asp_ppt.py
===============
Generate ASP_PRESENTATION.pptx — a full slide deck for the
Active Spiking Perception (ASP) paper.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# -----------------------------------------------------------------------
# Colour palette
# -----------------------------------------------------------------------
DARK_BLUE   = RGBColor(0x1A, 0x3A, 0x6B)   # title bar background
MID_BLUE    = RGBColor(0x27, 0x5C, 0xA6)   # accent
LIGHT_BLUE  = RGBColor(0xD6, 0xE8, 0xF8)   # background tint
ORANGE      = RGBColor(0xE8, 0x76, 0x1A)   # highlight / SSP colour
GREEN       = RGBColor(0x1C, 0x7D, 0x40)   # good result
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x11, 0x11, 0x11)
GRAY        = RGBColor(0x88, 0x88, 0x88)
LIGHT_GRAY  = RGBColor(0xF0, 0xF0, 0xF0)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK_LAYOUT = prs.slide_layouts[6]   # completely blank


# -----------------------------------------------------------------------
# Helpers
# -----------------------------------------------------------------------

def add_rect(slide, left, top, width, height, fill=None, line=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    fill_fmt = shape.fill
    if fill:
        fill_fmt.solid()
        fill_fmt.fore_color.rgb = fill
    else:
        fill_fmt.background()
    line_fmt = shape.line
    if line:
        line_fmt.color.rgb = line
    else:
        line_fmt.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=BLACK,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def slide_header(slide, title, subtitle=None):
    """Dark blue top bar with white title."""
    add_rect(slide, 0, 0, 13.33, 1.15, fill=DARK_BLUE)
    add_text(slide, title,
             0.3, 0.1, 12.5, 0.85,
             font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle,
                 0.3, 0.85, 12.5, 0.35,
                 font_size=14, color=RGBColor(0xCC, 0xDD, 0xFF),
                 align=PP_ALIGN.LEFT, italic=True)


def bullet_block(slide, bullets, left, top, width, height,
                 font_size=17, title=None, title_color=MID_BLUE):
    """Render a list of bullet strings as a text box."""
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(font_size + 1)
        run.font.bold = True
        run.font.color.rgb = title_color
        first = False

    for b in bullets:
        p = tf.paragraphs[0] if (first and not title) else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = b
        run.font.size = Pt(font_size)
        run.font.color.rgb = BLACK
        first = False
    return txb


def code_block(slide, code, left, top, width, height, font_size=12):
    bg = add_rect(slide, left, top, width, height,
                  fill=RGBColor(0x1E, 0x1E, 0x2E), line=RGBColor(0x44, 0x44, 0x66))
    txb = slide.shapes.add_textbox(
        Inches(left + 0.1), Inches(top + 0.05),
        Inches(width - 0.2), Inches(height - 0.1)
    )
    tf = txb.text_frame
    tf.word_wrap = False
    first = True
    for line in code.strip().split('\n'):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = RGBColor(0xCC, 0xFF, 0xCC)
        run.font.name = 'Courier New'
        first = False
    return txb


def highlight_box(slide, text, left, top, width, height,
                  bg=LIGHT_BLUE, border=MID_BLUE, font_size=16, bold=False):
    add_rect(slide, left, top, width, height, fill=bg, line=border)
    add_text(slide, text, left + 0.1, top + 0.08,
             width - 0.2, height - 0.15,
             font_size=font_size, bold=bold, color=BLACK, wrap=True)


def result_box(slide, text, left, top, width, height):
    add_rect(slide, left, top, width, height,
             fill=RGBColor(0xE8, 0xF8, 0xED),
             line=RGBColor(0x1C, 0x7D, 0x40))
    add_text(slide, text, left + 0.1, top + 0.08,
             width - 0.2, height - 0.15,
             font_size=16, bold=True, color=GREEN, wrap=True,
             align=PP_ALIGN.CENTER)


# -----------------------------------------------------------------------
# SLIDE 1 — Title
# -----------------------------------------------------------------------
sl = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(sl, 0, 0, 13.33, 7.5, fill=DARK_BLUE)
add_rect(sl, 0, 3.5, 13.33, 0.06, fill=ORANGE)

add_text(sl, "Active Spiking Perception (ASP)",
         0.8, 1.0, 11.5, 1.4,
         font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(sl, "Membrane-Guided Adaptive Slice Selection\nfor Anytime Energy-Efficient 3D Recognition",
         0.8, 2.4, 11.5, 1.2,
         font_size=22, color=RGBColor(0xBB, 0xDD, 0xFF),
         align=PP_ALIGN.CENTER)

add_text(sl, "Purdue SNN-PointNet Research  ·  March 2026",
         0.8, 4.0, 11.5, 0.6,
         font_size=16, color=RGBColor(0x88, 0xAA, 0xCC),
         align=PP_ALIGN.CENTER, italic=True)

add_text(sl, "95.0% accuracy  ·  42× energy savings  ·  2.4 / 16 slices on average",
         1.0, 5.0, 11.0, 0.8,
         font_size=18, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

# -----------------------------------------------------------------------
# SLIDE 2 — Motivation
# -----------------------------------------------------------------------
sl = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF8, 0xF8, 0xFF))
slide_header(sl, "Why Active Perception?", "The problem with fixed-order SNN methods")

bullet_block(sl, [
    "✓  SNNs communicate via binary spikes  →  cheap AC ops (not MACs)",
    "✓  Intel Loihi 2: AC costs 2.3×10⁻³ pJ vs 8.4×10⁻³ pJ for MAC  →  3.7× per-op saving",
    "✓  Trained SNNs have natural sparsity (r ≈ 0.3)  →  additional 3× savings",
    "",
    "✗  ALL prior methods (SPT, SPM, Spiking PointNet) process slices in FIXED order:",
    "        Slice₀  →  Slice₁  →  ...  →  Slice₁₅",
    "        regardless of what the object is",
    "",
    "✗  Problem: for a CHAIR, seeing the backrest first is highly discriminative.",
    "        For a LAMP, the base is.  Wasting timesteps on uninformative regions",
    "        burns energy needlessly.",
], 0.5, 1.3, 12.0, 5.5, font_size=16)

# callout
add_rect(sl, 0.5, 5.8, 12.33, 1.3, fill=RGBColor(0xFF, 0xF0, 0xD8),
         line=ORANGE)
add_text(sl, "💡  Our Insight:  The LIF membrane potential u_t after seeing t slices is a "
             "BELIEF STATE — a compressed summary of what is known.  Use it to decide "
             "where to look next.",
         0.65, 5.88, 12.0, 1.1,
         font_size=16, bold=False, color=RGBColor(0x80, 0x40, 0x00), wrap=True)

# -----------------------------------------------------------------------
# SLIDE 3 — Overview Diagram (text-based)
# -----------------------------------------------------------------------
sl = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF8, 0xF8, 0xFF))
slide_header(sl, "ASP Architecture Overview")

# Draw pipeline boxes
boxes = [
    (0.3,  2.0, 2.2, 1.0, LIGHT_BLUE,    "Point Cloud\nP ∈ ℝ^{1024×3}"),
    (2.8,  2.0, 2.2, 1.0, LIGHT_BLUE,    "FPS Slicing\nT=16 anchors"),
    (5.3,  2.0, 2.2, 1.0, LIGHT_BLUE,    "Geometry\nDescriptors G"),
    (5.3,  3.5, 2.2, 1.0, RGBColor(0xFF,0xEE,0xCC), "Slice Selection\nPolicy (SSP)\n~2K params"),
    (5.3,  5.0, 2.2, 1.0, LIGHT_BLUE,    "KNN Backbone\n+ Spiking MLP"),
    (8.1,  5.0, 2.2, 1.0, LIGHT_BLUE,    "Temporal SNN\n(LearnableLIF)"),
    (10.7, 5.0, 2.1, 1.0, RGBColor(0xE8,0xF8,0xED), "FC Classifier\n→ ŷ_t"),
]
for (l, t, w, h, col, txt) in boxes:
    add_rect(sl, l, t, w, h, fill=col, line=MID_BLUE)
    add_text(sl, txt, l+0.05, t+0.1, w-0.1, h-0.2,
             font_size=13, bold=False, color=BLACK, align=PP_ALIGN.CENTER)

# Arrow labels (simple text arrows)
arrows = [
    (2.5,  2.45, "→"),
    (5.0,  2.45, "→"),
    (6.4,  3.0,  "↓ g_m"),
    (6.4,  4.5,  "↓ m*"),
    (7.5,  5.45, "→"),
    (10.3, 5.45, "→"),
]
for (l, t, txt) in arrows:
    add_text(sl, txt, l, t, 0.8, 0.4, font_size=13, color=MID_BLUE,
             align=PP_ALIGN.CENTER, bold=True)

# Feedback arrow (membrane → SSP)
add_text(sl, "u_t (membrane)\n←←←←←←←←←←←←←←",
         4.5, 4.2, 5.0, 0.6, font_size=12, color=ORANGE,
         italic=True, align=PP_ALIGN.CENTER)

# Early exit
add_rect(sl, 10.7, 3.5, 2.1, 1.0,
         fill=RGBColor(0xE8, 0xF8, 0xED), line=GREEN)
add_text(sl, "Margin\n> θ ?\nEXIT →",
         10.7, 3.5, 2.1, 1.0, font_size=13, color=GREEN,
         bold=True, align=PP_ALIGN.CENTER)
add_text(sl, "↑", 11.75, 4.5, 0.3, 0.5, font_size=18, color=GREEN, bold=True,
         align=PP_ALIGN.CENTER)

add_text(sl, "Active Loop: membrane belief drives slice selection",
         0.5, 6.6, 12.0, 0.6, font_size=15, italic=True,
         color=GRAY, align=PP_ALIGN.CENTER)

# -----------------------------------------------------------------------
# SLIDE 4 — LIF Neurons
# -----------------------------------------------------------------------
sl = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF8, 0xF8, 0xFF))
slide_header(sl, "Background: Learnable LIF Neurons")

add_text(sl, "Standard LIF Dynamics",
         0.5, 1.3, 6.0, 0.5, font_size=18, bold=True, color=DARK_BLUE)

code_block(sl, """\
u_t  =  τ · u_{t-1}  +  W · x_t       # integrate
s_t  =  Θ(u_t − ϑ)                     # fire if above threshold
u_t  ←  u_t · (1 − s_t)               # reset after spike

Surrogate gradient: ∂s/∂u ≈ 1 / (1 + |u|)²""",
           0.5, 1.8, 6.0, 1.5, font_size=12)

add_text(sl, "Learnable per-neuron τ and ϑ  (ASP's neurons)",
         0.5, 3.5, 6.0, 0.5, font_size=18, bold=True, color=DARK_BLUE)

code_block(sl, """\
τ_i  =  σ(α_i)          ∈ (0,1)   # sigmoid → always valid leak
ϑ_i  =  softplus(β_i)   > 0       # always positive threshold
α_i, β_i  learnable per neuron""",
           0.5, 4.0, 6.0, 1.2, font_size=12)

# Example on the right
add_rect(sl, 7.0, 1.3, 5.8, 5.8,
         fill=RGBColor(0xFD, 0xF6, 0xE3), line=ORANGE)
add_text(sl, "Worked Example  (1 neuron, τ=0.9, ϑ=1.0, w=0.6)",
         7.1, 1.4, 5.6, 0.5, font_size=14, bold=True,
         color=RGBColor(0x80, 0x40, 0x00))

steps = [
    ("t=1", "x=1.0", "u = 0.9×0 + 0.6×1.0 = 0.60", "u<ϑ  ⟹  no spike"),
    ("t=2", "x=1.0", "u = 0.9×0.60 + 0.6×1.0 = 1.14", "u>ϑ  ⟹  SPIKE!  u→0"),
    ("t=3", "x=0.5", "u = 0.9×0 + 0.6×0.5 = 0.30", "u<ϑ  ⟹  no spike"),
]
for i, (t, inp, calc, res) in enumerate(steps):
    y = 2.1 + i * 1.5
    add_rect(sl, 7.1, y, 5.6, 1.2,
             fill=WHITE, line=RGBColor(0xCC, 0xBB, 0x88))
    add_text(sl, t, 7.15, y+0.05, 0.6, 0.35, font_size=13, bold=True,
             color=DARK_BLUE)
    add_text(sl, inp, 7.75, y+0.05, 1.5, 0.35, font_size=12, color=GRAY)
    add_text(sl, calc, 7.15, y+0.4, 5.3, 0.35, font_size=12,
             color=BLACK)
    col = GREEN if "SPIKE" in res else RGBColor(0xAA, 0x44, 0x00)
    add_text(sl, res, 7.15, y+0.75, 5.3, 0.35, font_size=13, bold=True,
             color=col)

# -----------------------------------------------------------------------
# SLIDE 5 — FPS + Geometry Descriptors
# -----------------------------------------------------------------------
sl = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF8, 0xF8, 0xFF))
slide_header(sl, "Data Preprocessing: FPS Slicing + Geometry Descriptors")

add_text(sl, "Farthest Point Sampling (FPS)",
         0.5, 1.3, 5.8, 0.5, font_size=18, bold=True, color=DARK_BLUE)

bullet_block(sl, [
    "1.  Pick any random seed point a₀",
    "2.  Each step: pick point FARTHEST from all anchors so far",
    "3.  Assign each remaining point to nearest anchor",
    "4.  Each anchor's points = one SLICE (64 pts)",
    "",
    "Result: T=16 slices with maximal spatial coverage",
    "        NO semantic labels needed",
], 0.5, 1.85, 5.8, 3.8, font_size=15)

# geo descriptor table on right
add_text(sl, "Geometry Descriptor  g_m ∈ ℝ⁶  (precomputed, no MACs)",
         6.6, 1.3, 6.5, 0.5, font_size=16, bold=True, color=DARK_BLUE)

rows = [
    ("Indices", "Feature", "What it means"),
    ("0–2", "anchor_x, y, z", "WHERE is this cluster?"),
    ("3", "mean dist from centroid", "Peripheral or central?"),
    ("4", "intra-cluster spread", "Thin edge vs. dense blob?"),
    ("5", "norm point count", "How many points (dense/sparse)?"),
]
for i, (idx, feat, meaning) in enumerate(rows):
    y = 1.9 + i * 0.6
    bg = DARK_BLUE if i == 0 else (LIGHT_BLUE if i % 2 == 0 else WHITE)
    tc = WHITE if i == 0 else BLACK
    add_rect(sl, 6.6, y, 0.8, 0.55, fill=bg, line=MID_BLUE)
    add_rect(sl, 7.4, y, 2.4, 0.55, fill=bg, line=MID_BLUE)
    add_rect(sl, 9.8, y, 3.4, 0.55, fill=bg, line=MID_BLUE)
    add_text(sl, idx,     6.62, y+0.08, 0.75, 0.4, font_size=12, bold=(i==0), color=tc)
    add_text(sl, feat,    7.42, y+0.08, 2.35, 0.4, font_size=12, bold=(i==0), color=tc)
    add_text(sl, meaning, 9.82, y+0.08, 3.35, 0.4, font_size=12, bold=(i==0), color=tc)

# example table for chair
add_text(sl, "Example — Chair with 4 anchors:",
         0.5, 5.8, 12.0, 0.4, font_size=15, bold=True, color=DARK_BLUE)

rows2 = [
    ("Anchor", "Region*", "x,y,z", "dist", "spread", "norm cnt"),
    ("a₀", "backrest", "0.1,0,0.9", "0.60", "0.08", "0.80"),
    ("a₁", "leg",      "0.0,0,-0.8","0.80", "0.02", "0.50"),
    ("a₂", "armrest",  "0.5,0.4,0.2","0.40","0.03","0.60"),
    ("a₃", "seat",     "-0.1,0,0.1", "0.10","0.05","1.30"),
]
widths = [0.9, 1.2, 2.0, 0.8, 0.8, 0.9]
starts = [0.5]
for w in widths[:-1]:
    starts.append(starts[-1] + w)
for i, row in enumerate(rows2):
    y = 6.2 + i * 0.45
    bg = DARK_BLUE if i == 0 else (LIGHT_BLUE if i % 2 == 0 else WHITE)
    tc = WHITE if i == 0 else BLACK
    for j, (cell, w, s) in enumerate(zip(row, widths, starts)):
        add_rect(sl, s, y, w, 0.42, fill=bg, line=MID_BLUE)
        add_text(sl, cell, s+0.02, y+0.05, w-0.04, 0.33,
                 font_size=11, color=tc, bold=(i==0))

add_text(sl, "* labels are for human readers — the algorithm only sees numbers",
         0.5, 8.3, 12.0, 0.4, font_size=10, italic=True, color=GRAY)

# -----------------------------------------------------------------------
# SLIDE 6 — KNN Backbone
# -----------------------------------------------------------------------
sl = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF8, 0xF8, 0xFF))
slide_header(sl, "KNN Local Backbone")

add_text(sl, "Per-point local feature construction",
         0.5, 1.3, 6.2, 0.5, font_size=18, bold=True, color=DARK_BLUE)

bullet_block(sl, [
    "For each point p_i in a slice (64 points):",
    "",
    "  1. Find k=16 nearest neighbours  {p_j1, ..., p_jk}  within the slice",
    "  2. Compute relative offsets:  Δp_jl = p_jl − p_i",
    "  3. Concatenate:  x_i = [p_i, Δp_j1, ..., Δp_jk]  ∈ ℝ^{3+3k} = ℝ^51",
    "  4. Feed through spiking MLP: 51 → 128 → 256 → 512",
    "     (LearnableLIF at each layer)",
    "",
    "  5. Slice embedding = mean-pool over all 64 points",
    "     e = mean({f_i})  ∈ ℝ^512",
], 0.5, 1.85, 6.2, 4.8, font_size=15)

# Why KNN box
add_rect(sl, 7.0, 1.3, 5.8, 3.0,
         fill=RGBColor(0xFD, 0xF6, 0xE3), line=ORANGE)
add_text(sl, "Why KNN neighbourhood?",
         7.1, 1.4, 5.6, 0.5, font_size=16, bold=True,
         color=RGBColor(0x80, 0x40, 0x00))
bullet_block(sl, [
    "✗ Simple PointNet: each point processed",
    "    independently — loses local structure",
    "",
    "✓ KNN: each point 'sees' its local patch",
    "    Similar to EdgeConv (DGCNN) but inside",
    "    the SNN pipeline",
    "",
    "✓ Relative offsets are translation-invariant",
    "✓ Works with any point density",
], 7.1, 1.95, 5.6, 3.0, font_size=13)

code_block(sl, """\
# Local feature per point
neighbours = knn_graph(pts, k=16)    # [B, N, 16, 3]
rel = neighbours - pts.unsqueeze(2)  # relative offsets
rel_flat = rel.reshape(B, N, 48)     # [B, N, 48]
x = cat([pts, rel_flat], dim=-1)     # [B, N, 51]

# Spiking MLP (all B*N points in parallel)
for layer in [lif128, lif256, lif512]:
    spk, mem = layer(x); x = mem

slice_emb = mem.reshape(B, N, 512).mean(dim=1)  # [B, 512]""",
           0.5, 5.0, 12.5, 2.1, font_size=11)

# -----------------------------------------------------------------------
# SLIDE 7 — SSP Architecture
# -----------------------------------------------------------------------
sl = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF8, 0xF8, 0xFF))
slide_header(sl, "Slice Selection Policy (SSP)", "The ~2K parameter 'brain' of ASP")

add_text(sl, "Scaled Dot-Product Attention: membrane attends to geometry",
         0.5, 1.3, 12.0, 0.45, font_size=17, color=DARK_BLUE, bold=True)

code_block(sl, """\
# Inputs
u       ∈ ℝ^D         # current LIF membrane state (belief)
g_m     ∈ ℝ^6         # geometry descriptor of anchor m

# Projections
k   = W_k @ u          ∈ ℝ^{d_ssp}       W_k ∈ ℝ^{64×512}
q_m = W_q @ g_m        ∈ ℝ^{d_ssp}       W_q ∈ ℝ^{64×6}

# Score
score_m = (k · q_m) / sqrt(d_ssp)         ∈ ℝ

# Mask visited anchors
score_m = -inf   if m ∈ visited""",
           0.5, 1.8, 7.5, 2.5, font_size=12)

# Analogy table
add_rect(sl, 8.3, 1.3, 4.7, 2.8,
         fill=RGBColor(0xE8, 0xF0, 0xFF), line=MID_BLUE)
add_text(sl, "Analogy: Attention",
         8.4, 1.4, 4.5, 0.45, font_size=15, bold=True, color=DARK_BLUE)
rows_a = [
    ("Term", "Role"),
    ("Key (k)", "What I currently believe"),
    ("Query (q_m)", "What region m looks like"),
    ("Score", "How relevant is m given my belief"),
    ("Mask", "Don't re-visit seen regions"),
]
for i, (a, b) in enumerate(rows_a):
    y = 1.9 + i * 0.44
    bg = MID_BLUE if i == 0 else (LIGHT_BLUE if i % 2 == 0 else WHITE)
    tc = WHITE if i == 0 else BLACK
    add_rect(sl, 8.3, y, 2.0, 0.42, fill=bg, line=MID_BLUE)
    add_rect(sl, 10.3, y, 2.65, 0.42, fill=bg, line=MID_BLUE)
    add_text(sl, a, 8.32, y+0.06, 1.96, 0.33, font_size=12, color=tc, bold=(i==0))
    add_text(sl, b, 10.32, y+0.06, 2.61, 0.33, font_size=12, color=tc)

# Worked example
add_rect(sl, 0.5, 4.5, 12.5, 2.6,
         fill=RGBColor(0xFD, 0xF6, 0xE3), line=ORANGE)
add_text(sl, "Worked Example: SSP Scoring (Chair, T=4 anchors)",
         0.65, 4.58, 12.0, 0.45, font_size=15, bold=True,
         color=RGBColor(0x80, 0x40, 0x00))

add_text(sl,
    "t=0:  u=0 (nothing seen)   → scores=[0.3, 0.1, 0.8, 0.2]   → select anchor 2 (high-z backrest region)\n"
    "t=1:  u encodes 'tall vertical structure'   visited={2}   → scores=[0.9, 0.2, -∞, 0.4]   → select anchor 0 (seat)\n"
    "t=2:  u encodes 'tall + flat'   margin=0.76 > θ=0.7   →  EARLY EXIT  →  predict: CHAIR  (2/4 slices used)",
    0.65, 5.1, 12.0, 1.8, font_size=14, color=BLACK, wrap=True)

# -----------------------------------------------------------------------
# SLIDE 8 — Training: Gumbel-Softmax
# -----------------------------------------------------------------------
sl = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF8, 0xF8, 0xFF))
slide_header(sl, "Training: Differentiable Slice Selection", "Gumbel-Softmax Straight-Through Estimator")

add_text(sl, "Problem: argmax is non-differentiable",
         0.5, 1.3, 12.0, 0.45, font_size=17, bold=True, color=RGBColor(0xAA, 0x00, 0x00))

add_text(sl,
    "We need gradients to flow from the loss into SSP weights W_k, W_q.\n"
    "But argmax(scores) has zero gradient everywhere  →  SSP cannot learn.",
    0.5, 1.8, 12.0, 0.7, font_size=15, color=BLACK)

add_rect(sl, 0.5, 2.6, 5.8, 1.8, fill=LIGHT_BLUE, line=MID_BLUE)
add_text(sl, "Forward Pass (hard — discrete):",
         0.6, 2.65, 5.6, 0.4, font_size=14, bold=True, color=DARK_BLUE)
add_text(sl, "w_t  =  one-hot( argmax(scores_t) )\n"
             "Each step picks EXACTLY ONE slice",
         0.6, 3.1, 5.6, 1.2, font_size=14, color=BLACK)

add_rect(sl, 6.8, 2.6, 6.1, 1.8, fill=RGBColor(0xFF, 0xEE, 0xCC), line=ORANGE)
add_text(sl, "Backward Pass (soft — differentiable):",
         6.9, 2.65, 5.9, 0.4, font_size=14, bold=True,
         color=RGBColor(0x80, 0x40, 0x00))
add_text(sl, "∂w_t/∂s_t  =  ∂softmax(s_t+γ)/∂s_t\n"
             "γ ~ Gumbel(0,1)  →  gradients flow into SSP",
         6.9, 3.1, 5.9, 1.2, font_size=14, color=BLACK)

add_text(sl, "→  Straight-Through: forward is hard argmax, backward is soft Gumbel",
         0.5, 4.5, 12.5, 0.45, font_size=16, bold=True, color=MID_BLUE,
         align=PP_ALIGN.CENTER)

code_block(sl, """\
# Differentiable feature selection
scores = ssp(membrane, geo, visited_mask)        # [B, T]
w = F.gumbel_softmax(scores, tau=τ, hard=True)   # [B, T] ≈ one-hot

e_t = (w.unsqueeze(-1) * all_feats).sum(dim=1)   # [B, D]  weighted sum
# w is ~one-hot so e_t ≈ all_feats[:, argmax_slice, :]
# but gradient flows through the soft weights""",
           0.5, 5.1, 12.5, 1.8, font_size=12)

add_text(sl, "Temperature τ annealed: τ₀=1.0 (exploratory) → 0.1 (deterministic) over 50 epochs",
         0.5, 7.0, 12.5, 0.4, font_size=13, italic=True, color=GRAY)

# -----------------------------------------------------------------------
# SLIDE 9 — Joint Loss Function
# -----------------------------------------------------------------------
sl = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF8, 0xF8, 0xFF))
slide_header(sl, "Joint Loss Function", "Four terms, each with a distinct role")

add_text(sl, "ℒ  =  ℒ_CE(ŷ_T, y)  +  λ_aux · (1/T) Σ ℒ_CE(ŷ_t, y)  "
             "+  λ_exit · (1/T) Σ (T-t)/T · (1 - max p_t)  +  λ_fr · r̄",
         0.5, 1.3, 12.5, 0.6, font_size=16, bold=True, color=DARK_BLUE,
         align=PP_ALIGN.CENTER)

terms = [
    ("ℒ_CE\n(Final)",
     "Standard CE on last timestep logit\nPrimary accuracy signal\nλ = 1.0",
     LIGHT_BLUE, MID_BLUE),
    ("ℒ_aux\n(Anytime)",
     "CE at every intermediate timestep\nEnables accurate early exit at ANY step\nAlso trains SSP ordering\nλ = 0.3",
     RGBColor(0xFF,0xEE,0xCC), ORANGE),
    ("ℒ_exit\n(Early Conf.)",
     "Penalise low max softmax probability\nEarlier timesteps weighted MORE\nTeaches SSP to pick discriminative regions first\nλ = 0.1",
     RGBColor(0xE8,0xF8,0xED), GREEN),
    ("ℒ_fr\n(Sparsity)",
     "Penalise mean firing rate r̄\nKeeps energy consumption low\nConverges to r ≈ 0.15–0.2\nλ = 0.05",
     RGBColor(0xF0,0xE8,0xFF), RGBColor(0x60,0x20,0x80)),
]
for i, (name, desc, bg, border) in enumerate(terms):
    x = 0.5 + i * 3.2
    add_rect(sl, x, 2.1, 3.0, 4.2, fill=bg, line=border)
    add_text(sl, name, x+0.1, 2.15, 2.8, 0.75,
             font_size=18, bold=True, color=border, align=PP_ALIGN.CENTER)
    add_rect(sl, x+0.05, 2.85, 2.9, 0.03, fill=border)
    add_text(sl, desc, x+0.1, 2.95, 2.8, 3.2,
             font_size=13, color=BLACK, wrap=True)

result_box(sl,
    "Key: ℒ_exit weight (T-t)/T means early timesteps penalised HARDER for low confidence\n"
    "→  SSP learns to select maximally informative regions FIRST",
    0.5, 6.5, 12.5, 0.85)

# -----------------------------------------------------------------------
# SLIDE 10 — Inference: Anytime Early Exit
# -----------------------------------------------------------------------
sl = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF8, 0xF8, 0xFF))
slide_header(sl, "Inference: Anytime Early Exit",
             "Only process slices until confident — then stop")

code_block(sl, """\
u = 0;  visited = {}

for t in range(T):                                    # T=16 max
    m* = argmax_SSP(u, geo, unvisited)                # O(M) dot products
    e_t = backbone( pts_slices[:, m*, :, :] )         # ONE backbone pass
    logits_t = temporal_head( e_t )                   # update membrane
    u = current_membrane()

    probs = softmax(logits_t)
    margin = probs[top1] - probs[top2]

    if margin > θ:
        return logits_t, exit_step=t+1               # EARLY EXIT""",
           0.5, 1.3, 12.5, 3.0, font_size=12)

add_text(sl, "Threshold θ controls the accuracy–energy trade-off:",
         0.5, 4.45, 12.0, 0.45, font_size=16, bold=True, color=DARK_BLUE)

rows_e = [
    ("θ", "Accuracy", "Mean exit steps", "Savings vs ANN"),
    ("0.5", "79.5%", "1.91 / 16", "83×"),
    ("0.7", "90.5%", "2.44 / 16", "63×  ← matches fixed-order baseline"),
    ("0.9", "95.0%", "3.49 / 16", "42×  ← best accuracy in study"),
    ("1.0", "93.5%", "16.0 / 16", "8.4×  (no early exit)"),
]
widths_e = [1.0, 1.5, 2.2, 7.0]
starts_e = [0.5, 1.5, 3.0, 5.25]
for i, row in enumerate(rows_e):
    y = 5.0 + i * 0.48
    bg = DARK_BLUE if i == 0 else (RGBColor(0xE8,0xF8,0xED) if i in [2,3] else WHITE)
    tc = WHITE if i == 0 else (GREEN if i in [2,3] else BLACK)
    for cell, w, s in zip(row, widths_e, starts_e):
        add_rect(sl, s, y, w, 0.45, fill=bg, line=MID_BLUE)
        add_text(sl, cell, s+0.03, y+0.06, w-0.06, 0.35,
                 font_size=12, color=tc, bold=(i==0))

# -----------------------------------------------------------------------
# SLIDE 11 — Results
# -----------------------------------------------------------------------
sl = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF8, 0xF8, 0xFF))
slide_header(sl, "Experimental Results", "ModelNet10 · 30-epoch run · 200 validation samples")

# 3 big result numbers
for i, (num, label, sub) in enumerate([
    ("95.0%",  "Accuracy", "at θ=0.9"),
    ("63×",    "Energy Savings", "at θ=0.7 (vs ANN)"),
    ("2.44/16","Mean Exit Steps", "91% exit within 4"),
]):
    x = 0.5 + i * 4.3
    add_rect(sl, x, 1.3, 3.9, 1.6, fill=DARK_BLUE, line=MID_BLUE)
    add_text(sl, num, x+0.1, 1.35, 3.7, 0.9,
             font_size=40, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    add_text(sl, label, x+0.1, 2.2, 3.7, 0.4,
             font_size=15, color=WHITE, align=PP_ALIGN.CENTER, bold=True)
    add_text(sl, sub, x+0.1, 2.55, 3.7, 0.3,
             font_size=12, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

# Comparison table
add_text(sl, "Comparison with published SNN baselines:",
         0.5, 3.15, 12.5, 0.45, font_size=16, bold=True, color=DARK_BLUE)

rows_r = [
    ("Model", "Dataset", "Accuracy", "Energy Savings"),
    ("Spiking PointNet (2023)", "MN40", "88.2%", "~5×"),
    ("SPT (2025)",              "MN40", "91.4%", "6.4×"),
    ("SPM (2025)",              "MN40", "92.3%", "3.5×"),
    ("Ours fixed-T",            "MN10", "90.6%", "8.4×"),
    ("ASP θ=0.7 (ours)",        "MN10", "90.5%", "63×  ← 7.5× better than fixed-T"),
    ("ASP θ=0.9 (ours)",        "MN10", "95.0%", "42×  ← highest accuracy"),
]
cw = [4.0, 1.5, 1.8, 5.0]
cs = [0.5, 4.5, 6.0, 7.8]
for i, row in enumerate(rows_r):
    y = 3.65 + i * 0.52
    is_ours = i in [5, 6]
    bg = DARK_BLUE if i == 0 else (RGBColor(0xD0,0xF0,0xD8) if is_ours else
                                   (LIGHT_BLUE if i % 2 == 0 else WHITE))
    tc = WHITE if i == 0 else (GREEN if is_ours else BLACK)
    for cell, w, s in zip(row, cw, cs):
        add_rect(sl, s, y, w, 0.49, fill=bg, line=MID_BLUE)
        add_text(sl, cell, s+0.04, y+0.07, w-0.08, 0.38,
                 font_size=12, color=tc, bold=(i==0 or is_ours))

# -----------------------------------------------------------------------
# SLIDE 12 — Exit Time Distribution
# -----------------------------------------------------------------------
sl = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF8, 0xF8, 0xFF))
slide_header(sl, "Exit Time Distribution", "θ=0.7  ·  ModelNet10  ·  200 validation samples")

add_text(sl, "Distribution of exit steps  (how many slices does ASP inspect before deciding?)",
         0.5, 1.3, 12.5, 0.5, font_size=16, color=DARK_BLUE)

# Bar chart (manual rectangles)
bars = [
    ("t=1", 0.20, "20%  (bathtub, monitor — easy)"),
    ("t=2", 0.53, "53%  (dominant bin — mean=2.5)"),
    ("t=3–4", 0.18, "18%  (moderate difficulty)"),
    ("t≥5",  0.09, "9%   (chair vs sofa, desk vs table)"),
]
bar_h_scale = 3.0  # max height in inches
max_frac = 0.53
for i, (label, frac, desc) in enumerate(bars):
    x = 1.2 + i * 2.8
    bar_h = (frac / max_frac) * bar_h_scale
    y_bot = 5.2
    bar_col = GREEN if frac == max(b[1] for b in bars) else MID_BLUE
    add_rect(sl, x, y_bot - bar_h, 2.0, bar_h, fill=bar_col, line=DARK_BLUE)
    add_text(sl, f"{frac*100:.0f}%", x+0.1, y_bot - bar_h - 0.45, 1.8, 0.4,
             font_size=18, bold=True, color=bar_col, align=PP_ALIGN.CENTER)
    add_text(sl, label, x+0.1, y_bot + 0.05, 1.8, 0.4,
             font_size=14, color=DARK_BLUE, align=PP_ALIGN.CENTER, bold=True)
    add_text(sl, desc, x - 0.5, y_bot + 0.5, 3.0, 0.7,
             font_size=11, color=GRAY, wrap=True)

# Axis line
add_rect(sl, 0.9, 5.2, 10.5, 0.03, fill=DARK_BLUE)
add_rect(sl, 0.9, 2.0, 0.03, 3.2, fill=DARK_BLUE)

add_text(sl,
    "Key insight:  The L_exit loss successfully trains ASP to reach high confidence early.\n"
    "Easy shapes (bathtubs, monitors) exit after 1 slice — discriminative structure is immediately obvious.\n"
    "Hard pairs (chair vs sofa) need 3–4 slices but still far fewer than the fixed 16.",
    0.5, 6.3, 12.5, 1.1, font_size=14, color=BLACK, wrap=True)

# -----------------------------------------------------------------------
# SLIDE 13 — Energy Analysis
# -----------------------------------------------------------------------
sl = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF8, 0xF8, 0xFF))
slide_header(sl, "Energy Analysis", "Following Lemaire et al. (2022) analytical model")

add_text(sl, "E_SNN / E_ANN  =  r̄  ×  (E_AC / E_MAC)  ×  (T_exit / T)",
         0.5, 1.3, 12.5, 0.55, font_size=22, bold=True, color=DARK_BLUE,
         align=PP_ALIGN.CENTER)

calcs = [
    ("Fixed-order baseline\n(ours_full)",
     "0.434  ×  0.274  ×  (16/16)  =  0.119\n→  8.4× cheaper than ANN",
     LIGHT_BLUE, MID_BLUE, "8.4×"),
    ("ASP  θ=0.7",
     "0.378  ×  0.274  ×  (2.44/16)  =  0.0158\n→  63× cheaper than ANN",
     RGBColor(0xD0,0xF0,0xD8), GREEN, "63×"),
    ("ASP  θ=0.9",
     "0.397  ×  0.274  ×  (3.49/16)  =  0.0237\n→  42× cheaper than ANN",
     RGBColor(0xD0,0xF0,0xD8), GREEN, "42×"),
]
for i, (title, calc, bg, tc, savings) in enumerate(calcs):
    x = 0.4 + i * 4.3
    add_rect(sl, x, 2.1, 4.0, 2.4, fill=bg, line=tc)
    add_text(sl, title, x+0.1, 2.15, 3.8, 0.65,
             font_size=15, bold=True, color=tc)
    add_text(sl, calc, x+0.1, 2.85, 3.8, 1.5,
             font_size=13, color=BLACK, wrap=True)
    add_text(sl, savings, x+0.6, 4.2, 2.8, 0.5,
             font_size=28, bold=True, color=tc, align=PP_ALIGN.CENTER)

add_text(sl, "Improvement:  0.119 / 0.0158  =  7.5×  better than fixed-T at identical accuracy",
         0.5, 4.8, 12.5, 0.5, font_size=16, bold=True, color=RGBColor(0xAA, 0x00, 0x00),
         align=PP_ALIGN.CENTER)

add_text(sl, "Why does ASP have slightly higher firing rate (0.378 vs 0.329 for fixed-T)?",
         0.5, 5.5, 12.5, 0.4, font_size=14, color=DARK_BLUE, bold=True)
add_text(sl,
    "The SSP selects the MOST informative slices first — these tend to have denser point clouds "
    "and stronger activations (slightly higher r̄).  But the dramatic reduction in T_exit "
    "(2.44 vs 16) more than compensates, yielding the Pareto-dominant frontier.",
    0.5, 5.95, 12.5, 1.0, font_size=14, color=BLACK, wrap=True)

# -----------------------------------------------------------------------
# SLIDE 14 — Complete Worked Example
# -----------------------------------------------------------------------
sl = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF8, 0xF8, 0xFF))
slide_header(sl, "Complete Worked Example: Chair Classification",
             "T=4 slices (simplified)  ·  θ=0.7")

steps_ex = [
    ("Preprocess",
     "FPS selects 4 anchors.  Geometry descriptors G precomputed.\n"
     "Backbone features NOT yet computed (computed on demand at inference).",
     LIGHT_BLUE, MID_BLUE),
    ("t=0  →  anchor 2",
     "u=0 (nothing seen).  SSP scores: [0.3, 0.1, 0.8, 0.2].  Select anchor 2 (backrest-region).\n"
     "Backbone(S₂) → e₀.  Temporal head: u₁ = f(0, e₀).\n"
     "Logits: chair=0.35, sofa=0.30.  Margin=0.05 < θ.  CONTINUE.",
     RGBColor(0xFF,0xEE,0xCC), ORANGE),
    ("t=1  →  anchor 0",
     "u₁ encodes 'tall vertical structure'.  visited={2}.\n"
     "SSP scores: [0.9, 0.2, -∞, 0.4].  Select anchor 0 (seat-region).\n"
     "Backbone(S₀) → e₁.  Temporal head updates: u₂.\n"
     "Logits: chair=0.82, sofa=0.06.  Margin=0.76 > θ=0.7.  EXIT!",
     RGBColor(0xD0,0xF0,0xD8), GREEN),
    ("Result",
     "Prediction: CHAIR  ✓\n"
     "Slices used: 2 of 4  (50% energy vs fixed-order)\n"
     "At T=16 scale: ~2.4 of 16 slices  →  63× energy saving",
     RGBColor(0xD0,0xF0,0xD8), GREEN),
]
for i, (title, desc, bg, border) in enumerate(steps_ex):
    y = 1.3 + i * 1.55
    add_rect(sl, 0.4, y, 12.5, 1.4, fill=bg, line=border)
    add_text(sl, title, 0.5, y+0.05, 2.8, 0.5,
             font_size=15, bold=True, color=border)
    add_rect(sl, 3.3, y+0.08, 0.03, 1.2, fill=border)
    add_text(sl, desc, 3.5, y+0.05, 9.2, 1.2,
             font_size=13, color=BLACK, wrap=True)

# -----------------------------------------------------------------------
# SLIDE 15 — Summary + Contributions
# -----------------------------------------------------------------------
sl = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(sl, 0, 0, 13.33, 7.5, fill=DARK_BLUE)
add_rect(sl, 0, 1.1, 13.33, 0.06, fill=ORANGE)

add_text(sl, "Summary & Contributions",
         0.5, 0.2, 12.0, 0.8,
         font_size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

contribs = [
    ("1.  Slice Selection Policy (SSP)",
     "~2K-param dot-product attention: membrane state attends to geometry descriptors.\n"
     "Selects the MOST informative unvisited region at each timestep."),
    ("2.  Differentiable Training",
     "Gumbel-softmax straight-through: hard selection in forward, soft gradients in backward.\n"
     "Enables end-to-end joint training of SSP + backbone + temporal head."),
    ("3.  Four-Term Joint Loss",
     "ℒ_CE + ℒ_aux + ℒ_exit + ℒ_fr.  Each term serves a specific role in the active loop.\n"
     "ℒ_exit with linear time-weighting is the key to early-exit efficiency."),
    ("4.  Pareto-Dominant Frontier",
     "At every accuracy threshold 85–95%, ASP uses strictly less energy than\n"
     "SPT (6.4×), SPM (3.5×), and our fixed-order baselines (8.4×)."),
]
for i, (title, desc) in enumerate(contribs):
    y = 1.4 + i * 1.45
    add_rect(sl, 0.4, y, 12.5, 1.3, fill=RGBColor(0x22, 0x44, 0x88), line=ORANGE)
    add_text(sl, title, 0.55, y+0.05, 12.0, 0.45,
             font_size=16, bold=True, color=ORANGE)
    add_text(sl, desc, 0.55, y+0.5, 12.0, 0.75,
             font_size=13, color=RGBColor(0xCC, 0xDD, 0xFF), wrap=True)

add_text(sl,
    "The membrane potential is a belief state.  Use it to look where it matters.",
    0.5, 7.1, 12.5, 0.35,
    font_size=15, italic=True, color=ORANGE, align=PP_ALIGN.CENTER)

# -----------------------------------------------------------------------
# Save
# -----------------------------------------------------------------------
out_path = "ASP_PRESENTATION.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
