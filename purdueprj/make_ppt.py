# -*- coding: utf-8 -*-
"""
Build a PowerPoint presentation for the SNN-PointNet experiment report.
Run from the purdueprj/ directory.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from copy import deepcopy

# ── paths ─────────────────────────────────────────────────────────────────────
BASE   = os.path.dirname(os.path.abspath(__file__))
IMGS   = os.path.join(BASE, "results")
OUT    = os.path.join(BASE, "report_slides.pptx")

# ── brand colours ─────────────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x1A, 0x1A, 0x2E)   # very dark navy
MID_BG     = RGBColor(0x16, 0x21, 0x3E)   # dark blue
ACCENT1    = RGBColor(0x0F, 0x3E, 0x85)   # Purdue-ish blue
ACCENT2    = RGBColor(0xE9, 0x4F, 0x37)   # red-orange (SNN)
ACCENT3    = RGBColor(0x39, 0xA0, 0xED)   # light blue (ANN)
GOLD       = RGBColor(0xCF, 0xB5, 0x37)   # Purdue gold
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xCC, 0xCC, 0xCC)
GREEN      = RGBColor(0x2E, 0xCC, 0x71)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]   # completely blank layout


# ── helpers ───────────────────────────────────────────────────────────────────
def rgb(r, g, b): return RGBColor(r, g, b)

def add_rect(slide, left, top, width, height, fill_color, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape

def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def add_image(slide, path, left, top, width=None, height=None):
    if not os.path.exists(path):
        return
    kw = {}
    if width:  kw["width"]  = Inches(width)
    if height: kw["height"] = Inches(height)
    return slide.shapes.add_picture(path, Inches(left), Inches(top), **kw)

def slide_bg(slide, color=DARK_BG):
    add_rect(slide, 0, 0, 13.33, 7.5, color)

def header_bar(slide, color=ACCENT1):
    add_rect(slide, 0, 0, 13.33, 0.9, color)

def accent_line(slide, left=0.4, top=0.88, width=12.53, color=GOLD, h=0.04):
    add_rect(slide, left, top, width, h, color)

def slide_number(slide, n, total):
    add_text(slide, f"{n} / {total}", 12.0, 7.1, 1.2, 0.3,
             font_size=11, color=LIGHT_GREY, align=PP_ALIGN.RIGHT)

def bullet_box(slide, items, left, top, width, height,
               font_size=16, color=WHITE, bullet="•", spacing=0.42):
    for i, item in enumerate(items):
        add_text(slide, f"{bullet}  {item}",
                 left, top + i * spacing, width, 0.4,
                 font_size=font_size, color=color)

def kv_row(slide, label, value, left, top, lw=2.8, vw=3.5,
           lsize=14, vsize=14, vcolor=GOLD):
    add_text(slide, label, left, top, lw, 0.35,
             font_size=lsize, color=LIGHT_GREY, bold=False)
    add_text(slide, value, left + lw, top, vw, 0.35,
             font_size=vsize, color=vcolor, bold=True)

TOTAL_SLIDES = 14

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, DARK_BG)
add_rect(s, 0, 0, 13.33, 0.06, GOLD)            # top gold strip
add_rect(s, 0, 7.44, 13.33, 0.06, GOLD)         # bottom gold strip
add_rect(s, 0, 2.6, 13.33, 2.3, MID_BG)         # mid panel

add_text(s, "Spiking vs. Standard PointNet",
         0.5, 1.2, 12.33, 1.0,
         font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "3-D Point Cloud Classification with Temporal Slicing & Early Exit",
         0.5, 2.05, 12.33, 0.6,
         font_size=20, color=GOLD, align=PP_ALIGN.CENTER, italic=True)

add_text(s, "PointNet-SNN  vs.  PointNet-ANN  on  ModelNet10",
         0.5, 2.8, 12.33, 0.55,
         font_size=22, bold=True, color=ACCENT3, align=PP_ALIGN.CENTER)

add_text(s, "Purdue Project — SNN Experiment", 0.5, 3.5, 12.33, 0.5,
         font_size=15, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
add_text(s, "February 19, 2026", 0.5, 3.85, 12.33, 0.4,
         font_size=14, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

# stat boxes at bottom
for i, (label, val, col) in enumerate([
    ("695K", "Parameters", WHITE),
    ("10", "Epochs", WHITE),
    ("ModelNet10", "Dataset", WHITE),
    ("16", "Time Slices", WHITE),
]):
    x = 1.0 + i * 2.9
    add_rect(s, x, 4.7, 2.5, 1.2, ACCENT1)
    add_text(s, label, x, 4.78, 2.5, 0.55,
             font_size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(s, val,   x, 5.28, 2.5, 0.4,
             font_size=13, color=WHITE, align=PP_ALIGN.CENTER)

slide_number(s, 1, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Outline
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s)
add_text(s, "Outline", 0.4, 0.12, 8, 0.65,
         font_size=28, bold=True, color=WHITE)
accent_line(s)
slide_number(s, 2, TOTAL_SLIDES)

sections = [
    ("01", "Motivation & Research Questions"),
    ("02", "Dataset — ModelNet10 & Radial Slicing"),
    ("03", "Model Architectures (ANN vs. SNN)"),
    ("04", "Training Setup & Loss Function"),
    ("05", "Training Results"),
    ("06", "Test Inference — 4 Modes"),
    ("07", "Accuracy vs. Timestep"),
    ("08", "Early Exit Analysis"),
    ("09", "Confidence Growth & Threshold Tradeoff"),
    ("10", "Analysis, Limitations & Next Steps"),
    ("11", "Conclusion"),
]

for i, (num, title) in enumerate(sections):
    col = 0 if i < 6 else 6.9
    row = (i % 6)
    y = 1.1 + row * 0.9
    add_rect(s, col + 0.4, y, 0.65, 0.55, ACCENT1)
    add_text(s, num, col + 0.4, y + 0.05, 0.65, 0.45,
             font_size=15, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(s, title, col + 1.15, y + 0.08, 5.6, 0.4,
             font_size=15, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Motivation
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s)
add_text(s, "Motivation & Research Questions", 0.4, 0.12, 12, 0.65,
         font_size=26, bold=True, color=WHITE)
accent_line(s)
slide_number(s, 3, TOTAL_SLIDES)

add_rect(s, 0.4, 1.05, 5.9, 5.9, MID_BG)
add_text(s, "Why Spiking Neural Networks?", 0.6, 1.15, 5.5, 0.5,
         font_size=17, bold=True, color=GOLD)
motivations = [
    "Event-driven, sparse computation",
    "Biologically inspired temporal dynamics",
    "Energy-efficient on neuromorphic hardware",
    "Natural fit for sequential / streaming data",
    "Integrate-and-fire = running evidence accumulator",
]
bullet_box(s, motivations, 0.6, 1.75, 5.5, 3.5, font_size=15, spacing=0.52)

add_rect(s, 6.8, 1.05, 6.1, 5.9, MID_BG)
add_text(s, "Research Questions", 7.0, 1.15, 5.7, 0.5,
         font_size=17, bold=True, color=GOLD)
questions = [
    "Can SNN-PointNet match standard ANN accuracy\n   on ModelNet10 with the same training budget?",
    "Does membrane potential accumulation across\n   slices produce progressively better predictions?",
    "Does confidence grow monotonically, enabling\n   meaningful early exit at test time?",
]
for i, q in enumerate(questions):
    add_rect(s, 7.0, 1.82 + i * 1.65, 5.7, 1.45, ACCENT1)
    add_text(s, f"Q{i+1}", 7.1, 1.9 + i * 1.65, 0.5, 0.4,
             font_size=16, bold=True, color=GOLD)
    add_text(s, q, 7.55, 1.88 + i * 1.65, 5.0, 1.2,
             font_size=13, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Dataset & Slicing
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s)
add_text(s, "Dataset — ModelNet10 & Radial Temporal Slicing", 0.4, 0.12, 12, 0.65,
         font_size=24, bold=True, color=WHITE)
accent_line(s)
slide_number(s, 4, TOTAL_SLIDES)

# left — dataset info
add_rect(s, 0.4, 1.05, 5.5, 5.9, MID_BG)
add_text(s, "ModelNet10", 0.6, 1.15, 5.1, 0.45,
         font_size=18, bold=True, color=GOLD)

stats = [
    ("Train samples", "3,991"),
    ("Test  samples", "908"),
    ("Classes", "10"),
    ("Points per cloud (N)", "1,024"),
    ("Point features", "(x, y, z) only"),
    ("File format", ".npy"),
]
for i, (k, v) in enumerate(stats):
    kv_row(s, k, v, 0.6, 1.75 + i * 0.52, lw=2.6, vw=2.7, lsize=13, vsize=14)

add_text(s, "10 Classes:", 0.6, 5.0, 5.2, 0.35,
         font_size=13, color=LIGHT_GREY, italic=True)
classes = "bathtub · bed · chair · desk · dresser\nmonitor · night stand · sofa · table · toilet"
add_text(s, classes, 0.6, 5.3, 5.2, 0.7, font_size=12, color=WHITE)

# right — slicing algorithm
add_rect(s, 6.3, 1.05, 6.6, 5.9, MID_BG)
add_text(s, "Radial Temporal Slicing", 6.5, 1.15, 6.2, 0.45,
         font_size=18, bold=True, color=GOLD)

steps = [
    "Compute centroid  mu = mean(points)",
    "Compute radial distance  d_i = ||p_i - mu||",
    "Sort points inner -> outer by distance",
    "Split into T=16 slices  (64 pts / slice)",
    "Feed slices sequentially to the network",
]
for i, step in enumerate(steps):
    add_rect(s, 6.5, 1.72 + i * 0.82, 0.45, 0.45, ACCENT1)
    add_text(s, str(i+1), 6.5, 1.78 + i * 0.82, 0.45, 0.35,
             font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, step, 7.05, 1.76 + i * 0.82, 5.7, 0.4,
             font_size=13, color=WHITE)

add_rect(s, 6.5, 5.85, 6.2, 0.75, ACCENT1)
add_text(s, "Coarse-to-fine: inner slices = shape skeleton,"
            "  outer slices = surface detail",
         6.6, 5.92, 6.0, 0.6, font_size=12, color=GOLD, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Architecture
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s)
add_text(s, "Model Architectures — ANN vs. SNN", 0.4, 0.12, 12, 0.65,
         font_size=26, bold=True, color=WHITE)
accent_line(s)
slide_number(s, 5, TOTAL_SLIDES)

# shared pipeline diagram
pipe_items = [
    ("Slice  S_t\n[B, 64, 3]", 0.5),
    ("Backbone MLP\n(per-point)", 2.5),
    ("Mean Pool\n[B, 512]", 4.8),
    ("Temporal Head\n(stateful)", 7.1),
    ("Logits\n[B, 10]", 9.6),
]
for i, (label, x) in enumerate(pipe_items):
    add_rect(s, x, 1.05, 1.7, 1.0, ACCENT1)
    add_text(s, label, x, 1.13, 1.7, 0.85,
             font_size=12, bold=(i in [1,3]), color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(pipe_items) - 1:
        add_text(s, "->", x + 1.7, 1.35, 0.6, 0.4,
                 font_size=18, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

add_text(s, "Shared macro-architecture (695,562 parameters in both models)",
         0.5, 2.15, 12.33, 0.35,
         font_size=12, color=LIGHT_GREY, italic=True, align=PP_ALIGN.CENTER)

# two-column comparison
add_rect(s, 0.4, 2.6, 6.0, 4.55, MID_BG)
add_rect(s, 0.4, 2.6, 6.0, 0.5, ACCENT3)
add_text(s, "PointNet-ANN  (Baseline)", 0.6, 2.67, 5.7, 0.38,
         font_size=16, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)

ann_rows = [
    ("Backbone", "3 x Linear + ReLU", "3->128->256->512"),
    ("Temporal head", "2 x Linear + ReLU", "512->512->512->10"),
    ("State across slices", "Cumulative mean", "of slice embeddings"),
    ("Activation", "ReLU (standard)", "f(x) = max(0,x)"),
    ("Generalisation", "76.2% test accuracy", "(after 10 epochs)"),
]
for i, (k, v1, v2) in enumerate(ann_rows):
    bg = MID_BG if i % 2 == 0 else rgb(0x1E, 0x30, 0x50)
    add_rect(s, 0.4, 3.18 + i * 0.76, 6.0, 0.74, bg)
    add_text(s, k,  0.55, 3.24 + i * 0.76, 2.0, 0.35, font_size=13, color=LIGHT_GREY)
    add_text(s, v1, 2.55, 3.24 + i * 0.76, 2.2, 0.35, font_size=13, bold=True, color=ACCENT3)
    add_text(s, v2, 4.75, 3.24 + i * 0.76, 1.5, 0.35, font_size=11, color=LIGHT_GREY)

add_rect(s, 6.9, 2.6, 6.0, 4.55, MID_BG)
add_rect(s, 6.9, 2.6, 6.0, 0.5, ACCENT2)
add_text(s, "PointNet-SNN  (Proposed)", 7.1, 2.67, 5.7, 0.38,
         font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

snn_rows = [
    ("Backbone", "3 x LIF layer", "3->128->256->512"),
    ("Temporal head", "2 x LIF + FC", "512->512->512->10"),
    ("State across slices", "Membrane potential", "u_t = tau*u_(t-1)+Wx"),
    ("Activation", "Heaviside + surrogate", "grad: 1/(1+|u|)^2"),
    ("Generalisation", "58.3% test accuracy", "(after 10 epochs)"),
]
for i, (k, v1, v2) in enumerate(snn_rows):
    bg = MID_BG if i % 2 == 0 else rgb(0x1E, 0x30, 0x50)
    add_rect(s, 6.9, 3.18 + i * 0.76, 6.0, 0.74, bg)
    add_text(s, k,  7.05, 3.24 + i * 0.76, 2.0, 0.35, font_size=13, color=LIGHT_GREY)
    add_text(s, v1, 9.05, 3.24 + i * 0.76, 2.2, 0.35, font_size=13, bold=True, color=ACCENT2)
    add_text(s, v2, 11.25, 3.24 + i * 0.76, 1.5, 0.35, font_size=11, color=LIGHT_GREY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — LIF Neuron & Loss Function
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s)
add_text(s, "LIF Neuron Dynamics & Training Loss", 0.4, 0.12, 12, 0.65,
         font_size=26, bold=True, color=WHITE)
accent_line(s)
slide_number(s, 6, TOTAL_SLIDES)

# LIF box
add_rect(s, 0.4, 1.05, 6.0, 5.9, MID_BG)
add_text(s, "Leaky Integrate-and-Fire (LIF)", 0.6, 1.15, 5.6, 0.45,
         font_size=17, bold=True, color=GOLD)

lif_eqs = [
    ("Membrane update:", "u_t  =  tau * u_(t-1)  +  W * x_t"),
    ("Spike emission:",  "s_t  =  Theta(u_t - theta)"),
    ("Soft reset:",      "u_t  <-  u_t * (1 - s_t)"),
    ("Surrogate grad:",  "ds/du  ~  1 / (1 + |u|)^2"),
]
for i, (label, eq) in enumerate(lif_eqs):
    add_text(s, label, 0.6, 1.75 + i * 0.85, 2.2, 0.35,
             font_size=13, color=LIGHT_GREY, italic=True)
    add_rect(s, 0.6, 2.05 + i * 0.85, 5.6, 0.45, ACCENT1)
    add_text(s, eq, 0.75, 2.1 + i * 0.85, 5.3, 0.38,
             font_size=15, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

params = [
    ("tau (decay)", "0.9"),
    ("theta (threshold)", "1.0"),
    ("Reset type", "Soft (subtract)"),
    ("Gradient approx.", "Surrogate"),
]
add_text(s, "Parameters", 0.6, 5.2, 5.5, 0.35, font_size=14,
         bold=True, color=GOLD)
for i, (k, v) in enumerate(params):
    kv_row(s, k, v, 0.6, 5.58 + i * 0.35, lw=2.5, vw=2.5, lsize=12, vsize=13)

# Loss function box
add_rect(s, 6.9, 1.05, 6.0, 5.9, MID_BG)
add_text(s, "Training Loss Function", 7.1, 1.15, 5.6, 0.45,
         font_size=17, bold=True, color=GOLD)

add_text(s, "L  =  L_CE(y_hat_T, y)  +  lambda * SUM L_CE(y_hat_t, y)",
         7.1, 1.78, 5.7, 0.55,
         font_size=13, bold=True, color=WHITE)
add_rect(s, 6.9, 2.42, 6.0, 0.06, GOLD)

components = [
    ("L_CE(y_hat_T, y)", "Final-slice loss", ACCENT3,
     "Supervises the last timestep (T=16).\nDrives end-to-end accuracy."),
    ("lambda * SUM L_CE(y_hat_t, y)", "Auxiliary loss  (lambda=0.3)", ACCENT2,
     "Supervises ALL intermediate steps.\nForces every slice to carry class\ninformation -> enables early exit."),
]
for i, (term, name, col, desc) in enumerate(components):
    add_rect(s, 7.1, 2.58 + i * 2.1, 5.7, 1.85, rgb(0x0F, 0x2D, 0x5A))
    add_text(s, term, 7.25, 2.65 + i * 2.1, 5.4, 0.45,
             font_size=14, bold=True, color=col)
    add_text(s, name, 7.25, 3.02 + i * 2.1, 5.4, 0.3,
             font_size=12, italic=True, color=LIGHT_GREY)
    add_text(s, desc, 7.25, 3.28 + i * 2.1, 5.4, 0.9,
             font_size=12, color=WHITE)

add_rect(s, 6.9, 6.75, 6.0, 0.2, ACCENT1)
add_text(s, "Optimizer: AdamW   |   LR: 1e-3   |   Grad clip: 1.0   |   Batch: 16",
         6.9, 6.82, 6.0, 0.3, font_size=11, color=GOLD, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Training Results
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s)
add_text(s, "Training Results — 10 Epochs", 0.4, 0.12, 12, 0.65,
         font_size=26, bold=True, color=WHITE)
accent_line(s)
slide_number(s, 7, TOTAL_SLIDES)

add_image(s, os.path.join(IMGS, "training_curves.png"),
          0.35, 0.98, width=8.5)

# stats panel
add_rect(s, 9.2, 0.98, 3.75, 6.45, MID_BG)
add_text(s, "Key Numbers", 9.35, 1.08, 3.4, 0.4,
         font_size=15, bold=True, color=GOLD)

kv_data = [
    ("ANN  Epoch 1",  "46.5%", ACCENT3),
    ("ANN  Epoch 10", "83.2%", ACCENT3),
    ("SNN  Epoch 1",  "35.0%", ACCENT2),
    ("SNN  Epoch 10", "87.4%", ACCENT2),
]
for i, (k, v, col) in enumerate(kv_data):
    add_rect(s, 9.35, 1.6 + i * 0.75, 3.4, 0.62,
             rgb(0x0F, 0x2D, 0x5A))
    add_text(s, k, 9.45, 1.67 + i * 0.75, 2.0, 0.3,
             font_size=12, color=LIGHT_GREY)
    add_text(s, v, 11.25, 1.62 + i * 0.75, 1.3, 0.45,
             font_size=20, bold=True, color=col, align=PP_ALIGN.RIGHT)

add_rect(s, 9.35, 4.65, 3.4, 1.6, ACCENT1)
add_text(s, "SNN starts slower\n(membrane calibration)\nbut surpasses ANN\nby epoch 2",
         9.45, 4.72, 3.2, 1.45, font_size=12, color=WHITE)

add_text(s, "Both models converge\nstably with no spikes\nor instability.",
         9.35, 6.35, 3.4, 0.9, font_size=11, color=LIGHT_GREY, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Inference: 4 Modes
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s)
add_text(s, "Test Inference — Four Evaluation Modes", 0.4, 0.12, 12, 0.65,
         font_size=26, bold=True, color=WHITE)
accent_line(s)
slide_number(s, 8, TOTAL_SLIDES)

modes = [
    ("ANN + Full",   "76.21%", "---",    ACCENT3,
     "Single forward pass on all 1024 points.\nNo slicing. Serves as the main baseline."),
    ("ANN + Slice",  "76.21%", "7.25",   ACCENT3,
     "Cumulative mean of slice embeddings.\nMatches Full; exits at step 7 on average."),
    ("SNN + Full",   " 9.25%", "---",    rgb(0xFF, 0x88, 0x44),
     "Single pass, one membrane step.\nNOT the intended mode — near-random\n(temporal dynamics not engaged)."),
    ("SNN + Slice",  "58.26%", "8.82",   ACCENT2,
     "Sequential slices + persistent membrane.\nIntended SNN mode. Mean exit: step 8.82."),
]

for i, (name, acc, exit_s, col, desc) in enumerate(modes):
    x = 0.35 + i * 3.25
    add_rect(s, x, 1.05, 3.1, 6.3, MID_BG)
    add_rect(s, x, 1.05, 3.1, 0.5, col)
    add_text(s, name, x, 1.1, 3.1, 0.42,
             font_size=15, bold=True, color=DARK_BG if col == ACCENT3 else WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s, "Test Accuracy", x, 1.7, 3.1, 0.3,
             font_size=11, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
    add_text(s, acc, x, 1.95, 3.1, 0.7,
             font_size=34, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_rect(s, x + 0.15, 2.78, 2.8, 0.04, col)
    add_text(s, "Mean Exit Step", x, 2.95, 3.1, 0.3,
             font_size=11, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
    add_text(s, exit_s + " / 16", x, 3.2, 3.1, 0.5,
             font_size=20, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(s, desc, x + 0.1, 3.85, 2.9, 1.8,
             font_size=12, color=WHITE)

# SNN+Full note
add_rect(s, 6.6, 1.05, 6.7, 0.72, rgb(0x5A, 0x20, 0x10))
add_text(s, "Why is SNN+Full near-random (9.25%)?",
         6.8, 1.1, 6.3, 0.35, font_size=14, bold=True, color=GOLD)
add_text(s,
    "The SNN temporal head is designed to integrate across 16 slices, with membrane "
    "potential acting as a running accumulator. Feeding all 1024 points in ONE step "
    "gives only a single integration tick — the membrane never builds meaningful state. "
    "The correct comparison is SNN+Slice (58.3%).",
    6.8, 1.77, 6.3, 1.5, font_size=12, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Accuracy vs Timestep
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s)
add_text(s, "Accuracy vs. Timestep", 0.4, 0.12, 12, 0.65,
         font_size=26, bold=True, color=WHITE)
accent_line(s)
slide_number(s, 9, TOTAL_SLIDES)

add_image(s, os.path.join(IMGS, "accuracy_vs_timestep.png"),
          0.35, 0.98, width=8.5)

add_rect(s, 9.2, 0.98, 3.75, 6.45, MID_BG)
add_text(s, "Observations", 9.35, 1.08, 3.4, 0.4,
         font_size=15, bold=True, color=GOLD)
obs = [
    "ANN rises steeply,\nplateaus by slice 6-8",
    "SNN shows slower but\nsustained growth",
    "Both exceed baselines\nwell before T=16",
    "Early exit viable\nfrom slice 6 onward",
]
for i, o in enumerate(obs):
    add_rect(s, 9.35, 1.6 + i * 1.35, 3.4, 1.2, ACCENT1)
    add_text(s, o, 9.45, 1.67 + i * 1.35, 3.2, 1.1,
             font_size=13, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Early Exit Analysis (histogram + CDF)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s)
add_text(s, "Early Exit Analysis", 0.4, 0.12, 12, 0.65,
         font_size=26, bold=True, color=WHITE)
accent_line(s)
slide_number(s, 10, TOTAL_SLIDES)

add_image(s, os.path.join(IMGS, "exit_histogram_snn.png"),
          0.35, 0.98, width=6.3)
add_image(s, os.path.join(IMGS, "exit_cdf.png"),
          6.75, 0.98, width=6.3)

add_text(s, "SNN exit histogram  (theta = 0.80)",
         0.35, 5.45, 6.3, 0.35, font_size=12, color=LIGHT_GREY, italic=True,
         align=PP_ALIGN.CENTER)
add_text(s, "Cumulative exit distribution",
         6.75, 5.45, 6.3, 0.35, font_size=12, color=LIGHT_GREY, italic=True,
         align=PP_ALIGN.CENTER)

add_rect(s, 0.35, 5.88, 12.63, 1.48, MID_BG)
add_text(s,
    "ANN exits earlier on average (mean 7.25) vs. SNN (mean 8.82), "
    "reflecting faster confidence build-up via cumulative averaging. "
    "Both models allow substantial compute savings: at theta=0.80, roughly half the slices "
    "can be skipped on average without waiting for the final step.",
    0.5, 5.95, 12.3, 1.35, font_size=13, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Confidence & Threshold Tradeoff
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s)
add_text(s, "Confidence Growth & Threshold Tradeoff", 0.4, 0.12, 12, 0.65,
         font_size=24, bold=True, color=WHITE)
accent_line(s)
slide_number(s, 11, TOTAL_SLIDES)

add_image(s, os.path.join(IMGS, "confidence_growth.png"),
          0.35, 0.98, width=6.3)
add_image(s, os.path.join(IMGS, "threshold_tradeoff.png"),
          6.75, 0.98, width=6.3)

add_text(s, "Avg. max-softmax confidence per slice",
         0.35, 5.45, 6.3, 0.35, font_size=12, color=LIGHT_GREY, italic=True,
         align=PP_ALIGN.CENTER)
add_text(s, "Accuracy vs. mean exit step (theta sweep)",
         6.75, 5.45, 6.3, 0.35, font_size=12, color=LIGHT_GREY, italic=True,
         align=PP_ALIGN.CENTER)

add_rect(s, 0.35, 5.88, 12.63, 1.48, MID_BG)
add_text(s,
    "Confidence grows monotonically for both models — validating temporal slicing.   "
    "The threshold curve (right) sweeps theta in [0.50, 0.99]: low theta gives early exit "
    "with lower accuracy; high theta defers to the final slice and recovers full accuracy. "
    "This tradeoff can be tuned at inference time with no retraining.",
    0.5, 5.95, 12.3, 1.35, font_size=13, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Analysis & Limitations
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s)
add_text(s, "Analysis & Limitations", 0.4, 0.12, 12, 0.65,
         font_size=26, bold=True, color=WHITE)
accent_line(s)
slide_number(s, 12, TOTAL_SLIDES)

add_rect(s, 0.4, 1.05, 6.0, 5.9, MID_BG)
add_text(s, "Key Findings", 0.6, 1.15, 5.6, 0.45,
         font_size=17, bold=True, color=GREEN)
findings = [
    "SNN achieves higher training acc (87.4%)\nthan ANN (83.2%) in 10 epochs",
    "Test gap (ANN 76% vs SNN 58%) signals\noverfitting in the SNN",
    "Radial slicing enables early exit:\n~45% avg compute savings at theta=0.80",
    "Confidence grows monotonically for both\nmodels — early exit is viable",
    "SNN+Full is invalid for this architecture\n(temporal dynamics not engaged)",
]
for i, f in enumerate(findings):
    add_rect(s, 0.5, 1.75 + i * 0.98, 5.7, 0.85, rgb(0x0A, 0x3A, 0x2A))
    add_text(s, "✓", 0.6, 1.83 + i * 0.98, 0.4, 0.5,
             font_size=16, bold=True, color=GREEN)
    add_text(s, f, 0.95, 1.83 + i * 0.98, 5.1, 0.75,
             font_size=12, color=WHITE)

add_rect(s, 6.9, 1.05, 6.0, 5.9, MID_BG)
add_text(s, "Limitations", 7.1, 1.15, 5.6, 0.45,
         font_size=17, bold=True, color=ACCENT2)
limits = [
    ("No GPU", "Training on CPU only; 30+ epochs\nfeasible with GPU in same time"),
    ("Fixed LR", "No schedule used; cosine/step decay\nexpected to improve SNN generalisation"),
    ("Small dataset", "ModelNet10 has only 3991 train\nsamples; ModelNet40/ScanObjectNN\nwould stress-test more"),
    ("No augmentation", "Random rotation, jitter, scale\nwould reduce overfitting significantly"),
]
for i, (title, desc) in enumerate(limits):
    add_rect(s, 7.1, 1.75 + i * 1.22, 5.6, 1.08, rgb(0x3A, 0x10, 0x10))
    add_text(s, title, 7.2, 1.82 + i * 1.22, 5.3, 0.35,
             font_size=13, bold=True, color=ACCENT2)
    add_text(s, desc, 7.2, 2.12 + i * 1.22, 5.3, 0.65,
             font_size=11, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Next Steps
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
header_bar(s)
add_text(s, "Recommended Next Steps", 0.4, 0.12, 12, 0.65,
         font_size=26, bold=True, color=WHITE)
accent_line(s)
slide_number(s, 13, TOTAL_SLIDES)

steps_data = [
    ("LR Scheduling",        "Cosine decay or step schedule;\n30+ epoch training on GPU",       ACCENT3),
    ("Data Augmentation",    "Random rotation, jitter, scale;\nexpected to close train-test gap", GREEN),
    ("Regularisation",       "Dropout in temporal head;\nweight decay tuning for SNN",           ACCENT3),
    ("Learnable Thresholds", "Per-layer adaptive LIF thresholds;\nlearnable tau parameter",      GREEN),
    ("Larger Benchmarks",    "ModelNet40, ShapeNet, or\nScanObjectNN for broader coverage",     ACCENT3),
    ("Energy Profiling",     "Measure spike rates and SynapticOps\nto quantify energy savings",  GREEN),
]

for i, (title, desc, col) in enumerate(steps_data):
    col_idx = i % 3
    row_idx = i // 3
    x = 0.4 + col_idx * 4.3
    y = 1.1 + row_idx * 2.8
    add_rect(s, x, y, 4.0, 2.55, MID_BG)
    add_rect(s, x, y, 4.0, 0.52, col)
    add_text(s, f"{i+1:02d}  {title}", x + 0.1, y + 0.07, 3.8, 0.4,
             font_size=14, bold=True,
             color=DARK_BG if col == GREEN else WHITE)
    add_text(s, desc, x + 0.15, y + 0.65, 3.7, 1.75,
             font_size=13, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Conclusion
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s, DARK_BG)
add_rect(s, 0, 0, 13.33, 0.06, GOLD)
add_rect(s, 0, 7.44, 13.33, 0.06, GOLD)
add_rect(s, 0, 1.55, 13.33, 0.06, ACCENT1)
add_rect(s, 0, 6.15, 13.33, 0.06, ACCENT1)

add_text(s, "Conclusion", 0.5, 0.2, 12.33, 0.9,
         font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s,
    "We presented a Spiking PointNet that processes 3-D point clouds as a\n"
    "temporal sequence of radial slices, enabling biologically motivated\n"
    "integrate-and-fire dynamics for 3-D shape classification.",
    0.5, 1.72, 12.33, 1.1,
    font_size=16, color=WHITE, align=PP_ALIGN.CENTER, italic=True)

boxes = [
    ("SNN Train Acc", "87.4%", ACCENT2,
     "Surpasses ANN training\naccuracy (83.2%)"),
    ("ANN Test Acc",  "76.2%", ACCENT3,
     "Strong ANN baseline;\ngap to close with more training"),
    ("SNN Test Acc",  "58.3%", GOLD,
     "Slice mode only;\nroom to improve with regularisation"),
    ("Early Exit",    "~45%",  GREEN,
     "Avg compute savings\nat theta = 0.80"),
]
for i, (label, val, col, note) in enumerate(boxes):
    x = 0.5 + i * 3.1
    add_rect(s, x, 3.0, 2.8, 2.85, MID_BG)
    add_rect(s, x, 3.0, 2.8, 0.45, col)
    add_text(s, label, x, 3.05, 2.8, 0.38,
             font_size=12, bold=True,
             color=DARK_BG if col in (GREEN, GOLD, ACCENT3) else WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s, val, x, 3.52, 2.8, 0.8,
             font_size=30, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s, note, x, 4.3, 2.8, 1.4,
             font_size=12, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s,
    "These results demonstrate the viability of spiking temporal processing for "
    "point cloud classification and motivate further work on GPU-accelerated, "
    "regularised SNN training for 3-D perception.",
    0.5, 6.22, 12.33, 0.9,
    font_size=14, color=LIGHT_GREY, align=PP_ALIGN.CENTER, italic=True)

slide_number(s, 14, TOTAL_SLIDES)

# ── save ──────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {TOTAL_SLIDES}")
